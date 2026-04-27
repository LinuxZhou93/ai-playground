import docx
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor

def generate_docx():
    doc = docx.Document()
    
    # Title
    title = doc.add_heading('天府七中墨子实验室校园开放日（5.17）展示方案', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Background
    doc.add_heading('一、 活动背景与目标', level=1)
    p = doc.add_paragraph('5月17日是天府七中的校园开放日，预计接待2000-3000名新初一学生及家长。本次活动旨在通过高科技、强互动的展示，树立墨子实验室“材料科学+机器人”的科技领先形象，促进招生转化。')
    
    # Location
    doc.add_heading('二、 点位安排与氛围营造', level=1)
    p = doc.add_paragraph()
    p.add_run('核心展示区：').bold = True
    p.add_run('报告厅大门口（人流必经C位）。')
    
    p = doc.add_paragraph()
    p.add_run('视觉物料：').bold = True
    p.add_run('在该区域架设大型桁架广告牌（约6m x 3m），主题定为“碳基未来：墨子实验室探索之旅”。结合行架展示墨子实验室的科技成就与课程体系。')
    
    # Projects
    doc.add_heading('三、 互动展示项目（精选自2026产品手册）', level=1)
    
    items = [
        ("1. 众擎机器人巡游表演", "由众擎仿生机器人进行现场行走、避障及互动表演，作为吸引人流的“超级吸铁石”。"),
        ("2. Carbon-X 石墨烯电阻云图测试", "源自材料科学课程L01。学生使用万用表实地测试石墨烯纸的电阻分布，感受新型碳材料的导电魅力。"),
        ("3. 碳基极客工坊：电路点亮LED", "学生使用导电笔或石墨烯材料现场搭建简易电路，点亮LED，完成科技小制作并可带走（低成本高转化）。"),
        ("4. 航空航天材料模型展示", "配合天府七中宜金校区的氛围，展示微纳卫星或空间站材料模型，突出材料在尖端科技中的应用。")
    ]
    
    for title_text, desc in items:
        p = doc.add_paragraph()
        p.add_run(title_text).bold = True
        p.add_run(f'\n{desc}')
    
    # Workflow
    doc.add_heading('四、 流程安排', level=1)
    table = doc.add_table(rows=1, cols=3)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = '时间段'
    hdr_cells[1].text = '活动内容'
    hdr_cells[2].text = '负责人'
    
    times = [
        ('08:30 - 09:00', '行架搭建、设备调试', '工程部'),
        ('09:00 - 11:30', '上午场巡演及互动测试', '教学部/导员'),
        ('11:30 - 13:30', '设备维护、午间轮班', '值班员'),
        ('13:30 - 16:30', '下午场重点展示及家长咨询', '咨询部/佳宁老师')
    ]
    
    for time, content, staff in times:
        row_cells = table.add_row().cells
        row_cells[0].text = time
        row_cells[1].text = content
        row_cells[2].text = staff

    # Cost
    doc.add_heading('五、 成本控制建议', level=1)
    doc.add_paragraph('• 优先使用现有机器人及实验器材。\n• 消耗性物料（石墨烯纸、LED等）按3000份标准集采，控制单人体验成本在5元以内。\n• 广告牌设计沿用视觉识别系统，减少二次设计费用。')

    doc.save('/Users/zhoulin/Desktop/517天府七中墨子实验室开放日策划方案.docx')

def generate_pptx():
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "5.17 天府七中开放日展示方案"
    subtitle.text = "墨子实验室：碳基未来 · 智领苍穹\n策划：墨子团队"
    
    # Slide 2: Background
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "活动背景 & 目标"
    content = slide.placeholders[1]
    content.text = "• 目标群体：2000-3000名新初一学生及家长\n• 核心诉求：招生转化、科技实力展示\n• 核心理念：跨学科（材料+机器人）体验"
    
    # Slide 3: Visual & Layout
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "点位安排与视觉氛围"
    content = slide.placeholders[1]
    content.text = "• 位置：报告厅门口（核心人流区）\n• 物料：大型行架广告牌（6m x 3m）\n• 风格：高保真科技感，深色调配合霓虹发光色"
    
    # Slide 4: Highlights 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "展示核心：众擎机器人"
    content = slide.placeholders[1]
    content.text = "• 项目：众擎仿生机器人现场互动\n• 亮点：动态巡游、语音交互、避障展示\n• 作用：制造视觉焦点，引发朋友圈社交传播"
    
    # Slide 5: Highlights 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "课程精选：Carbon-X 材料实验"
    content = slide.placeholders[1]
    content.text = "• 实验1：电阻云图测试（感知石墨烯导电性）\n• 实验2：碳基笔绘制导电回路（DIY科技小制作）\n• 目标：强调材料科学的创新应用"
    
    # Slide 6: Schedule
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "流程与成本"
    content = slide.placeholders[1]
    content.text = "• 09:00-16:30：全天候互动展示\n• 配合宜金校区招生策略进行针对性宣讲\n• 严格控制单人耗材成本在5元以内"

    prs.save('/Users/zhoulin/Desktop/517天府七中墨子实验室开放日策划方案.pptx')

if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    print("Files generated successfully on Desktop.")
