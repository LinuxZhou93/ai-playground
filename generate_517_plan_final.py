import docx
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- THEME COLORS ---
MOZI_DARK = PptRGBColor(10, 10, 12)
MOZI_CYAN = PptRGBColor(0, 240, 255)
MOZI_YELLOW = PptRGBColor(255, 214, 0)
MOZI_WHITE = PptRGBColor(226, 232, 240)

def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_styled_title(slide, text, color=MOZI_CYAN):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.color.rgb = color
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = PptPt(44)

def generate_docx():
    doc = docx.Document()
    
    # Title
    title = doc.add_heading('天府七中墨子实验室校园开放日（5.17）全案策划方案', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Lab Info
    doc.add_heading('一、 实验室使命与愿景', level=1)
    p = doc.add_paragraph('墨子实验室（Mozi Lab）专注于“全栈科技特长生”的闭环培养。我们秉承“自我管理（Self-Management）”的核心教育理念，通过“竞技导向”与“工程导向”双线并行，助力学生在材料科学、人工智能及机器人工程领域实现跨越式发展。')
    
    # Background
    doc.add_heading('二、 5.17 活动核心策略', level=1)
    p = doc.add_paragraph('目标受众：2000-3000名新初一学者及家庭。\n展示逻辑：视觉冲击（众擎机器人）-> 深度交互（石墨烯云图）-> 品牌连接（扫码领样品）-> 长期留存（火星课伏笔）。')
    
    # Layout
    doc.add_heading('三、 展位布局与视觉系统', level=1)
    p = doc.add_paragraph('点位：报告厅门口（流量咽喉）。\n视觉：6m x 3m 科技蓝/荧光绿主调桁架。结合墨子学者培养图谱 V15.0 的核心要素。')
    
    # Projects
    doc.add_heading('四、 互动展示项目精选', level=1)
    items = [
        ("1. 众擎机器人：全地形仿生交互", "由众擎仿生机器人进行实时避障及语音交互展示，体现实验室在AI机器人领域的尖端教研实力。"),
        ("2. 石墨烯伴手礼：扫码领“未来”", "每位学生扫码加微信可领取石墨烯薄膜样品一套。卡片内含“墨子学者”成长路径图，为后期转化埋下强有力钩子。"),
        ("3. 火星/月球基地：未来课程前瞻", "展示“星际殖民”系列课程模型，通过第一性原理（First Principles）引导学生思考地外生存挑战。"),
        ("4. Carbon-X 电阻云图测试", "现场提供高精度万用表，让学生亲手测绘新材料电阻，感受从原子到工程的跨度。")
    ]
    for title_text, desc in items:
        p = doc.add_paragraph()
        run = p.add_run(title_text)
        run.bold = True
        p.add_run(f'\n{desc}')

    doc.save('/Users/zhoulin/Desktop/517天府七中墨子实验室策划方案_高保真版.docx')

def generate_pptx():
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    set_slide_background(slide, MOZI_DARK)
    
    # Add a decorative shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(3), PptInches(10), PptInches(0.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MOZI_CYAN
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "5.17 天府七中校园开放日"
    p.font.bold = True
    p.font.size = PptPt(54)
    p.font.color.rgb = MOZI_CYAN
    
    txBox2 = slide.shapes.add_textbox(PptInches(1), PptInches(4), PptInches(8), PptInches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "墨子实验室：全栈科技特长生培养展示全案\n碳基极客 · 智绘未来"
    p2.font.size = PptPt(24)
    p2.font.color.rgb = MOZI_WHITE

    # Slide 2: Mission
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, MOZI_DARK)
    add_styled_title(slide, "实验室基因：墨子学者")
    content = slide.placeholders[1]
    content.text = "• 核心理念：自我管理（Self-Management）\n• 培养路径：全栈科技特长生培养图谱 V15.0\n• 核心赛道：VEX机器人竞技 + 工程技术创客"
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = MOZI_WHITE
        p.font.size = PptPt(24)

    # Slide 3: Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, MOZI_DARK)
    add_styled_title(slide, "5.17 活动执行逻辑")
    content = slide.placeholders[1]
    content.text = "• 流量入口：众擎机器人动态巡游（视觉爆点）\n• 转化钩子：扫码领石墨烯样品（伴手礼引流）\n• 伏笔植入：火星基地探索课程（长期期待）"
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = MOZI_WHITE
        p.font.size = PptPt(24)

    # Slide 4: Project Focus 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, MOZI_DARK)
    add_styled_title(slide, "重点展示：Carbon-X 系列")
    content = slide.placeholders[1]
    content.text = "• 石墨烯电阻云图实测：科学探究思维显化\n• 碳基极客工坊：DIY 导电电路明信片\n• 意义：将抽象材料科学转化为可感知的互动"
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = MOZI_WHITE
        p.font.size = PptPt(24)

    # Slide 5: Project Focus 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, MOZI_DARK)
    add_styled_title(slide, "生态延伸：火星/月球基地", MOZI_YELLOW)
    content = slide.placeholders[1]
    content.text = "• 替代航空航天静态展示：发布未来基地建构课程\n• 核心：第一性原理引导下的太空生存挑战\n• 目标：吸引高年龄段（新初一）极客学员"
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = MOZI_WHITE
        p.font.size = PptPt(24)

    # Slide 6: Operation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, MOZI_DARK)
    add_styled_title(slide, "运营与闭环")
    content = slide.placeholders[1]
    content.text = "• 佳宁老师团队：负责现场留资与深度咨询\n• 成本控制：耗材人均 < 5元，复用高端器材\n• 目标转化：预约 6 月份实验室深度开放日"
    for p in content.text_frame.paragraphs:
        p.font.color.rgb = MOZI_WHITE
        p.font.size = PptPt(24)

    prs.save('/Users/zhoulin/Desktop/517天府七中墨子实验室策划方案_高保真版.pptx')

if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    print("Branded High-Fidelity files generated on Desktop.")
