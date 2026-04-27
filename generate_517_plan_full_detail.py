import docx
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# --- MOZI BRANDING ---
MOZI_DARK = PptRGBColor(7, 8, 12)
MOZI_CYAN = PptRGBColor(0, 240, 255)
MOZI_YELLOW = PptRGBColor(255, 214, 0)
MOZI_WHITE = PptRGBColor(245, 245, 247)
MOZI_DEEP_BLUE = PptRGBColor(15, 23, 42)

def apply_text_style(paragraph, font_size=12, bold=False, color=None):
    run = paragraph.add_run()
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return run

def add_full_plan_content(doc):
    # Title
    t = doc.add_heading('【5.17天府七中】墨子实验室校园开放日·全场景全流程深度执行全案', 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # 1. Vision
    doc.add_heading('一、 品牌战略与教育愿景 (Mission & Vision)', level=1)
    doc.add_paragraph('墨子实验室不仅仅是科创空间，它是“全栈科技特长生”的孵化器。本次开放日的目标是在2小时的接触时间内，完成从“科技震撼”到“品牌认同”再到“生源入库”的全链路转化。')
    p = doc.add_paragraph('核心教育哲学：\n1. 第一性原理思维：引导学生从物理本质解决工程问题，而非经验模仿。\n2. 六维学者评估：从工程实践、计算思维、团队协作、抗挫能力、学术严谨、工具管理六个维度定义未来的科创领袖。')

    # 2. Layout & Visuals
    doc.add_heading('二、 现场点位与视觉系统 (Site & Visual Identity)', level=1)
    doc.add_paragraph('展示区：报告厅正大门广场（人流绝对核心区）。\n桁架参数：6米(W) x 3.5米(H) 加厚喷绘，主KV采用“赛博极客”风格。\n多维空间感：设置3个独立的互动岛（交互岛、制造岛、咨询岛），形成顺时针导流轨迹。')

    # 3. Project Detail (3x Expansion)
    doc.add_heading('三、 互动展示项目详细执行手册 (Project Manual)', level=1)
    
    # Unitree
    doc.add_heading('3.1 众擎仿生机器人：全场景交互表演 (The Anchor)', level=2)
    doc.add_paragraph('表演脚本：\n- 0-5min: 自由行走、人群招呼，建立亲和力。\n- 5-15min: 台阶攀爬、草地越障（体现天府七中真实地形适配）。\n- 15-20min: 语音对答，植入实验室Slogan：“想和机器人一起改变世界吗？”\n技术保障：双电池组热替换，冷却风扇全时段待命。')

    # Carbon-X
    doc.add_heading('3.2 Carbon-X：石墨烯电阻云图实测 (The Core)', level=2)
    doc.add_paragraph('背景：源自石墨烯精英班核心课。 \n详细步骤：\n1. 领取：学生领取一份5x5坐标网格纸与万用表。\n2. 测绘：在老师引导下，测试石墨烯纸样5个核心坐标的电阻值。\n3. 云图：根据数据在iPad端或纸质卡片上连线，观察电流走势。\n互动词：苏格拉底式引导——“为什么这边的电阻比那边大？如果我剪开一个小口，电流会怎么走？”')

    # Mars
    doc.add_heading('3.3 伏笔项目：火星基地探索课程 (The Seed)', level=2)
    doc.add_paragraph('展示内容：火星液氧甲烷火箭模型、地外生存密封仓教具。\n策略：展示但不完全开放体验，仅开放给意向报名的学员预约“首航体验”。\n物料：课程导览手册（包含北京大学相关课程延伸背景）。')

    # 4. Conversion Funnel
    doc.add_heading('四、 流量获取与转化 SOP (Conversion Funnel)', level=1)
    doc.add_paragraph('1. A级流量点：桁架前机器人拍照，导流至体验区。\n2. B级留存点：扫码领取“石墨烯样品礼包”。（礼包内含：样品、导电笔抵扣券、实验室门票）。\n3. C级转化点：佳宁老师团队现场完成《科技特长生潜质评估初表》录入。')

    # 5. Logistics & Risk
    doc.add_heading('五、 后勤保障与风险管理 (Logistics & Risk)', level=1)
    table = doc.add_table(rows=1, cols=3)
    table.style = 'Table Grid'
    hdr = table.rows[0].cells
    hdr[0].text = '潜在风险'
    hdr[1].text = '影响等级'
    hdr[2].text = '应急响应方案'
    
    risks = [
        ('机器人电力耗尽', '高', '准备4组高能锂电，单组满电巡航45min，轮流热替换。'),
        ('客流拥挤导致秩序混乱', '中', '设置单向参观流线，分批次放行（每批15人）。'),
        ('网络信号差无法扫码', '中', '提前打印纸质预约码，提供离线留资登记表。'),
        ('设备雨天防护', '高', '配备3x6米防水帐篷及临时防水罩。')
    ]
    for r, l, s in risks:
        row = table.add_row().cells
        row[0].text = r
        row[1].text = l
        row[2].text = s

    # 6. Budget
    doc.add_heading('六、 预算概算 (Budget Overview)', level=1)
    doc.add_paragraph('物料：大型桁架(2k) + 伴手礼(1.5k/300份) + 实验耗材(0.8k) = 4.3k\n人力：复用墨子团队自有人员。')

def generate_docx():
    doc = docx.Document()
    add_full_plan_content(doc)
    doc.save('/Users/zhoulin/Desktop/517天府七中开放日全案_全量细节版.docx')

def generate_pptx():
    prs = Presentation()
    
    # --- SLIDE 1: COVER (High Design) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    # Shapes
    add_color_block(slide, 0, 0, 10, 0.2, MOZI_CYAN)
    add_color_block(slide, 0, 0.2, 0.2, 7.5, MOZI_CYAN)
    
    tx = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(2))
    p = tx.text_frame.paragraphs[0]
    p.text = "5.17 天府七中开放日"
    p.font.size = PptPt(64)
    p.font.bold = True
    p.font.color.rgb = MOZI_CYAN
    
    p2 = tx.text_frame.add_paragraph()
    p2.text = "墨子实验室：深度引流与展示全案"
    p2.font.size = PptPt(32)
    p2.font.color.rgb = MOZI_WHITE
    
    # --- SLIDE 2: MISSION (Two Column) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "01 | 品牌主张：培养科技特长生")
    # Columns
    add_color_block(slide, 0.5, 2, 4.2, 4, MOZI_DEEP_BLUE)
    add_color_block(slide, 5.3, 2, 4.2, 4, MOZI_DEEP_BLUE)
    
    tx1 = slide.shapes.add_textbox(PptInches(0.6), PptInches(2.2), PptInches(4), PptInches(3))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "【核心使命】\n基于墨子学者培养图谱\nV15.0 体系，通过跨学科\n项目驱动自驱力。"
    p1.font.size = PptPt(18)
    p1.font.color.rgb = MOZI_WHITE
    
    tx2 = slide.shapes.add_textbox(PptInches(5.4), PptInches(2.2), PptInches(4), PptInches(3))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "【展示逻辑】\n从感官震撼到理性探究\n众擎机器人(A) -> \n石墨烯实验(B) -> \n留资转化(C)"
    p2.font.size = PptPt(18)
    p2.font.color.rgb = MOZI_WHITE

    # --- SLIDE 3: FUNNEL (Visual Shapes) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "02 | 流量漏斗与转化闭环")
    
    colors = [MOZI_CYAN, MOZI_CYAN, MOZI_CYAN, MOZI_YELLOW]
    steps = ["触达：机器人演练", "留存：石墨烯实验", "线索：样品礼包领取", "转化：佳宁团队咨询"]
    for i in range(4):
        w = 8 - i*1.5
        x = (10 - w)/2
        y = 2 + i*1.1
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i]
        shape.text_frame.text = steps[i]
        shape.text_frame.paragraphs[0].font.size = PptPt(16)
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].font.color.rgb = MOZI_DARK

    # --- SLIDE 4: ROBOT (Performance Script) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "03 | 众擎机器人：巡游脚本")
    add_color_block(slide, 0.5, 1.5, 9, 5, MOZI_DEEP_BLUE)
    tx = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.8), PptInches(8.5), PptInches(4.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• 极客表演：每逢整点进行 15 分钟“挑战赛”。\n• 互动点：学生可指令机器人进行“握手”、“搜寻”。\n• 传播：设置“未来指挥官”拍照合影牌。\n• 后勤：双组电池热切换 + 专业导员讲解。"
    for para in tf.paragraphs:
        para.font.size = PptPt(22)
        para.font.color.rgb = MOZI_WHITE

    # --- SLIDE 5: CARBON-X (Detailed Steps) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "04 | Carbon-X：石墨烯深度实验")
    add_color_block(slide, 0.5, 1.5, 9, 5, MOZI_DEEP_BLUE)
    tx = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.8), PptInches(8.5), PptInches(4.5))
    tf = tx.text_frame
    tf.text = "• 实验：电阻云图实测（全方位感知材料性能）。\n• 引导：苏格拉底式提问，建立探究闭环。\n• 产出：绘制专属“电流星图”并打卡留念。\n• 转化：凭实验记录单免费领取石墨烯样品。"
    for para in tf.paragraphs:
        para.font.size = PptPt(22)
        para.font.color.rgb = MOZI_WHITE

    # --- SLIDE 6: MARS BASE (Future Seed) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "05 | 伏笔：火星基地探索课程", MOZI_YELLOW)
    add_color_block(slide, 0.5, 1.5, 9, 5, MOZI_DEEP_BLUE)
    tx = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.8), PptInches(8.5), PptInches(4.5))
    tf = tx.text_frame
    tf.text = "• 展示：地外生存密封仓 + 能源采集教具。\n• 逻辑：强调“不仅仅是模型，更是工程挑战”。\n• 引流：仅开放 50 名“火星先锋官”试听席位。\n• 目的：将今日流量转化为秋季班核心存量。"
    for para in tf.paragraphs:
        para.font.size = PptPt(22)
        para.font.color.rgb = MOZI_WHITE

    # --- SLIDE 7: TEAM & RISK ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "06 | 团队协作与风险管控")
    add_color_block(slide, 0.5, 1.5, 4.3, 5, MOZI_DEEP_BLUE)
    add_color_block(slide, 5.2, 1.5, 4.3, 5, MOZI_DEEP_BLUE)
    
    tx1 = slide.shapes.add_textbox(PptInches(0.6), PptInches(1.7), PptInches(4), PptInches(4))
    tx1.text_frame.text = "【分工】\n• 教学：项目演示\n• 咨询：家长转化\n• 工程：搭建维保"
    for p in tx1.text_frame.paragraphs: p.font.size = PptPt(20); p.font.color.rgb = MOZI_WHITE

    tx2 = slide.shapes.add_textbox(PptInches(5.3), PptInches(1.7), PptInches(4), PptInches(4))
    tx2.text_frame.text = "【预案】\n• 机器人：电力双备\n• 客流：S型围栏导流\n• 网络：4G离线应急"
    for p in tx2.text_frame.paragraphs: p.font.size = PptPt(20); p.font.color.rgb = MOZI_WHITE

    prs.save('/Users/zhoulin/Desktop/517天府七中开放日全案_全量细节版.pptx')

def add_manual_title(slide, text, color=MOZI_CYAN):
    txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(9), PptInches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = color
    p.font.bold = True
    p.font.size = PptPt(36)

def add_color_block(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    print("Full-Detail Branded files generated on Desktop.")
