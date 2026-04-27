import docx
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- MOZI BRANDING ---
MOZI_DARK = PptRGBColor(7, 8, 12)
MOZI_CYAN = PptRGBColor(0, 240, 255)
MOZI_YELLOW = PptRGBColor(255, 214, 0)
MOZI_WHITE = PptRGBColor(245, 245, 247)
MOZI_DEEP_BLUE = PptRGBColor(15, 23, 42)

# --- HELPERS ---
def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_manual_title(slide, text, color=MOZI_CYAN):
    txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(9), PptInches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = color
    p.font.bold = True
    p.font.size = PptPt(36)

def add_color_block(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# --- MAIN FUNCTIONS ---
def generate_docx():
    doc = docx.Document()
    t = doc.add_heading('【5.17天府七中】墨子实验室校园开放日·全场景全流程深度执行全案', 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_heading('一、 品牌战略与教育愿景', level=1)
    doc.add_paragraph('墨子实验室不仅是科技展示场，更是未来科创人才的孵化器。本次活动旨在通过“极致视觉”与“深度探究”建立品牌护城河。核心理念包括“第一性原理”与“自我管理”。')
    
    doc.add_heading('二、 全量文字策划方案（细节扩充 3x）', level=1)
    
    sections = [
        ("2.1 现场布防系统", "桁架：采用 6m x 3.5m 异形定制桁架，主KV采用霓虹青与暗黑工业风结合。设置“极客入口”闸机感装饰。\n交互岛：配置 4 台高性能移动工作站，实时展示机器人运动算法模型。"),
        ("2.2 众擎机器人巡游详细脚本", "巡游路径：从报告厅正门出发，绕广场一周。重点在台阶处停留，展示动态平衡能力。\n语音交互：内置 50 条关于材料科学与天府七中特色的应答话术。例如：“我身上的涂层就是你们等会要领取的石墨烯材料哦”。\n应急：现场配备两名随行导员，手持遥控器进行“安全熔断”操作。"),
        ("2.3 Carbon-X 实验深度 SOP", "实验物料：300 份单原子层石墨烯纸片、纳米导电笔、高精度万用表。\n步骤：1. 观察微观形貌图；2. 进行电阻多点采样；3. 绘制等位线；4. 领取成品样品。"),
        ("2.4 佳宁老师团队：咨询转化闭环", "留资激励：扫码进入“天府七中极客交流群”，即可获得《石墨烯与未来教育》白皮书电子版。内容埋点：白皮书第 5 页详细介绍了秋季“墨子学者”精英班的选拔流程。")
    ]
    for st, sd in sections:
        p = doc.add_paragraph()
        p.add_run(st).bold = True
        p.add_run(f'\n{sd}')
    
    doc.add_heading('三、 风险管理手册', level=1)
    doc.add_paragraph('设备风险：众擎机器人配备 3 套冗余动力系统。耗材风险：石墨烯样品溢价备货 50%。客流风险：采取分时段“派票”进入互动区模式。')
    
    doc.save('/Users/zhoulin/Desktop/517天府七中开放日全案_全量细节版.docx')

def generate_pptx():
    prs = Presentation()
    
    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_color_block(slide, 0, 0, 10, 0.2, MOZI_CYAN)
    add_color_block(slide, 0, 0.2, 0.2, 7.5, MOZI_CYAN)
    tx = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(2))
    p = tx.text_frame.paragraphs[0]
    p.text = "5.17 天府七中开放日"
    p.font.size = PptPt(64); p.font.bold = True; p.font.color.rgb = MOZI_CYAN
    p2 = tx.text_frame.add_paragraph()
    p2.text = "墨子实验室：全量深度执行全案"; p2.font.size = PptPt(32); p2.font.color.rgb = MOZI_WHITE

    # Slide 2: Mission
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "01 | 品牌核心：全栈特长生培养")
    add_color_block(slide, 0.5, 1.5, 9, 5, MOZI_DEEP_BLUE)
    tx = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.8), PptInches(8.5), PptInches(4.5))
    tf = tx.text_frame
    tf.text = "• 愿景：通过极致科技项目驱动人才自进化\n• 体系：基于 V15.0 图谱的六维评估体系\n• 目标：在天七建立墨子实验室的“极客高地”"
    for p in tf.paragraphs: p.font.size = PptPt(24); p.font.color.rgb = MOZI_WHITE

    # Slide 3: Projects
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "02 | 项目精选：众擎机器人 & Carbon-X")
    add_color_block(slide, 0.5, 1.5, 4.3, 5, MOZI_DEEP_BLUE)
    add_color_block(slide, 5.2, 1.5, 4.3, 5, MOZI_DEEP_BLUE)
    tx1 = slide.shapes.add_textbox(PptInches(0.6), PptInches(1.8), PptInches(4), PptInches(4))
    tx1.text_frame.text = "【众擎巡演】\n全地形避障巡航\n实时语音对答互动\n合影拍照流量裂变"
    for p in tx1.text_frame.paragraphs: p.font.size = PptPt(20); p.font.color.rgb = MOZI_WHITE
    tx2 = slide.shapes.add_textbox(PptInches(5.3), PptInches(1.8), PptInches(4), PptInches(4))
    tx2.text_frame.text = "【石墨烯实验】\n电阻云图实测探究\nDIY 导电电路明信片\n伴手礼扫码精准转化"
    for p in tx2.text_frame.paragraphs: p.font.size = PptPt(20); p.font.color.rgb = MOZI_WHITE

    prs.save('/Users/zhoulin/Desktop/517天府七中开放日全案_全量细节版.pptx')

if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    print("Full-Detail Branded files generated on Desktop.")
