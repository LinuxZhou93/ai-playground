import docx
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- GOOGLE STYLE BRANDING ---
G_BLUE = PptRGBColor(66, 133, 244)
G_RED = PptRGBColor(219, 68, 55)
G_YELLOW = PptRGBColor(244, 180, 0)
G_GREEN = PptRGBColor(15, 157, 88)
G_WHITE = PptRGBColor(255, 255, 255)
G_BG = PptRGBColor(248, 249, 250)
G_TEXT = PptRGBColor(60, 64, 67)
G_GRAY = PptRGBColor(95, 99, 104)

def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_google_title(slide, text, color=G_BLUE):
    txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(9), PptInches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = color
    p.font.bold = True
    p.font.size = PptPt(32)

def add_rounded_block(slide, x, y, w, h, color, text="", text_color=G_WHITE, font_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        shape.text_frame.text = text
        p = shape.text_frame.paragraphs[0]
        p.font.color.rgb = text_color
        p.font.size = PptPt(font_size)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    return shape

def generate_docx():
    doc = docx.Document()
    t = doc.add_heading('天府七中 5.17 开放日：墨子实验室展示方案 (Google 简约科技风)', 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_heading('一、 项目背景与使命', level=1)
    doc.add_paragraph('墨子实验室致力于通过极客教育培养未来的科创领袖。本次活动采用 Google 简约设计风格，强调透明、高效、且具有亲和力的科技互动体验。')
    
    doc.add_heading('二、 详细展示脚本 (全量版)', level=1)
    p = doc.add_paragraph('【众擎机器人】\n表演逻辑：1. 自动巡航；2. 避障演示；3. 实时交互。每场表演时长15分钟，配以专业导员解说。\n【Carbon-X 实验】\n步骤：1. 分发石墨烯薄膜；2. 指导多点位采样；3. 扫码录入数据生成专属报告。\n【招生意向闭环】\n通过扫码领取“石墨烯样品礼包”，将人流引向企业微信，由佳宁老师团队进行后续深度咨询。')
    
    doc.save('/Users/zhoulin/Desktop/517天府七中开放日全案_谷歌风版.docx')

def generate_pptx():
    prs = Presentation()
    
    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, G_WHITE)
    
    # Google Accent Bar
    add_rounded_block(slide, 1, 1.5, 0.2, 3, G_BLUE)
    add_rounded_block(slide, 1.3, 1.5, 0.2, 3, G_RED)
    add_rounded_block(slide, 1.6, 1.5, 0.2, 3, G_YELLOW)
    add_rounded_block(slide, 1.9, 1.5, 0.2, 3, G_GREEN)
    
    tx = slide.shapes.add_textbox(PptInches(2.5), PptInches(1.5), PptInches(6), PptInches(3))
    p = tx.text_frame.paragraphs[0]
    p.text = "5.17 天府七中开放日"
    p.font.size = PptPt(48); p.font.bold = True; p.font.color.rgb = G_TEXT
    p2 = tx.text_frame.add_paragraph()
    p2.text = "墨子实验室：深度展示与引流全案\nGoogle Design Style"; p2.font.size = PptPt(24); p2.font.color.rgb = G_GRAY

    # Slide 2: Mission
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, G_BG)
    add_google_title(slide, "活动核心：培养精英学者")
    
    # 3 Cards
    add_rounded_block(slide, 0.5, 1.5, 2.8, 4, G_WHITE)
    add_rounded_block(slide, 3.6, 1.5, 2.8, 4, G_WHITE)
    add_rounded_block(slide, 6.7, 1.5, 2.8, 4, G_WHITE)
    
    for i, title in enumerate(["核心理念", "培养体系", "目标群体"]):
        tx = slide.shapes.add_textbox(PptInches(0.6 + i*3.1), PptInches(1.7), PptInches(2.6), PptInches(3.5))
        tf = tx.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title; p.font.bold = True; p.font.color.rgb = G_BLUE
        p2 = tf.add_paragraph()
        p2.text = ["自我管理\n第一性原理", "V15.0 培养图谱\n双向发展路径", "新初一准极客\n科创特长生"][i]
        p2.font.size = PptPt(16); p2.font.color.rgb = G_TEXT

    # Slide 3: Funnel
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, G_BG)
    add_google_title(slide, "转化漏斗：从流量到生源")
    
    colors = [G_BLUE, G_RED, G_YELLOW, G_GREEN]
    steps = ["触达", "互动", "线索", "转化"]
    descs = ["众擎机器人巡演", "石墨烯电阻测试", "扫码领样品礼包", "佳宁团队咨询"]
    for i in range(4):
        add_rounded_block(slide, 0.5 + i*2.35, 2.5, 2.2, 2.5, colors[i], text=steps[i])
        tx = slide.shapes.add_textbox(PptInches(0.5 + i*2.35), PptInches(5.2), PptInches(2.2), PptInches(1))
        p = tx.text_frame.paragraphs[0]
        p.text = descs[i]; p.font.size = PptPt(12); p.alignment = PP_ALIGN.CENTER; p.font.color.rgb = G_TEXT

    # Slide 4: Robot
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, G_WHITE)
    add_google_title(slide, "众擎机器人：巡游脚本")
    add_rounded_block(slide, 0.5, 1.5, 9, 5, G_BG)
    tx = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.8), PptInches(8.5), PptInches(4.5))
    tf = tx.text_frame
    tf.text = "• 极客表演：全地形避障与语音交互\n• 视觉焦点：双机巡游，打造震撼入场感\n• 传播互动：设置未来指挥官打卡区\n• 后勤保障：480min 巡航，双电热切备选"
    for p in tf.paragraphs: p.font.size = PptPt(22); p.font.color.rgb = G_TEXT

    # Slide 5: Carbon-X
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, G_WHITE)
    add_google_title(slide, "Carbon-X：深度实验互动", G_RED)
    add_rounded_block(slide, 0.5, 1.5, 9, 5, G_BG)
    tx = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.8), PptInches(8.5), PptInches(4.5))
    tf = tx.text_frame
    tf.text = "• 实验：单层石墨烯电阻测绘云图\n• 引导：苏格拉底式提问，建立探究思维\n• 闭环：扫码领取石墨烯样品礼包\n• 转化：由咨询部完成潜质评估档案录入"
    for p in tf.paragraphs: p.font.size = PptPt(22); p.font.color.rgb = G_TEXT

    prs.save('/Users/zhoulin/Desktop/517天府七中开放日全案_谷歌风版.pptx')

if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    print("Google-Style Branded files generated on Desktop.")
