import docx
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- THEME COLORS ---
MOZI_DARK = PptRGBColor(10, 10, 12)
MOZI_CYAN = PptRGBColor(0, 240, 255)
MOZI_YELLOW = PptRGBColor(255, 214, 0)
MOZI_WHITE = PptRGBColor(226, 232, 240)
MOZI_GRAY = PptRGBColor(100, 116, 139)

def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_manual_title(slide, text, color=MOZI_CYAN):
    txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(9), PptInches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = color
    p.font.bold = True
    p.font.size = PptPt(36)

def add_color_block(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def generate_docx():
    doc = docx.Document()
    
    # Title
    title = doc.add_heading('天府七中墨子实验室 5.17 开放日展示全案策划（深度执行版）', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 1. 核心使命
    doc.add_heading('一、 核心使命与教育哲学', level=1)
    p = doc.add_paragraph('墨子实验室（Mozi Lab）作为全栈科技特长生培养基地，本次开放日的核心不仅是产品展示，更是“精英学者”价值观的输出。我们强调：\n• 第一性原理（First Principles）：培养从物理本质思考问题的深度。\n• 自我管理（Self-Management）：通过高难度科创项目驱动学生自主成长。')
    
    # 2. 活动漏斗 (Mermaid)
    doc.add_heading('二、 活动转化漏斗 (User Journey Flow)', level=1)
    doc.add_paragraph('以下为本次活动的流量转化逻辑（Mermaid）：')
    p = doc.add_paragraph('graph TD\nA[大型行架+众擎机器人巡游] -->|吸引停留| B[石墨烯电阻云图互动]\nB -->|深度兴趣| C[扫码领石墨烯样品/留资]\nC -->|意向挖掘| D[佳宁老师团队专业咨询]\nD -->|长期伏笔| E[预约实验室深度开放日/暑期营]')
    
    # 3. 详细项目执行计划
    doc.add_heading('三、 详细项目执行计划', level=1)
    
    items = [
        ("1. 众擎机器人巡游表演", "规格：由2台众擎仿生机器人组成巡游方阵。\n脚本：每30分钟进行一次整点表演（避障、上下楼梯、语音应答）。\n目标：制造全场视觉高地，引导“自媒体”二次传播。"),
        ("2. Carbon-X：电阻云图实测", "实验步骤：学生使用万用表在5x5坐标网上测试石墨烯纸电阻值，记录数据并绘制等势线。\n核心：理解材料非均质性，感受从实验室到工程化的过程。"),
        ("3. 探索课程引流：火星基地/月球能源", "展示内容：Mars-Alpha基地模型、地外资源采矿教具。\n伏笔逻辑：介绍这不仅是模型，而是即将开启的“星际移民”系列挑战赛课，吸引准初一学生。"),
        ("4. 科技节产品精选：导电电路 DIY", "物料：导电胶带、3V纽扣电池、LED发光二极管、定制明信片。\n产出：学生亲手制作完成并发光的“墨子学者”明信片。")
    ]
    for t, d in items:
        p = doc.add_paragraph()
        run = p.add_run(t)
        run.bold = True
        p.add_run(f'\n{d}')

    # 4. 团队分工
    doc.add_heading('四、 团队权责清单', level=1)
    doc.add_paragraph('• 教学部：负责机器人演示、实验引导、苏格拉底式提问。\n• 咨询部（佳宁老师）：负责意向家长沟通、留资表填写、后期回访。\n• 工程部：负责桁架搭建、设备电力保障、现场安全。')

    doc.save('/Users/zhoulin/Desktop/517天府七中墨子实验室策划方案_终极详细版.docx')

def generate_pptx():
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_color_block(slide, 0, 0, 10, 0.5, MOZI_CYAN)
    add_color_block(slide, 0, 7, 10, 0.5, MOZI_YELLOW)
    
    txBox = slide.shapes.add_textbox(PptInches(1), PptInches(2), PptInches(8), PptInches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "天府七中 5.17 开放日"
    p.font.bold = True
    p.font.size = PptPt(54)
    p.font.color.rgb = MOZI_CYAN
    p2 = tf.add_paragraph()
    p2.text = "墨子实验室深度策划全案"
    p2.font.size = PptPt(36)
    p2.font.color.rgb = MOZI_WHITE

    # Slide 2: Mission
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "核心使命：培养全栈科技特长生")
    add_color_block(slide, 0.5, 1.2, 0.1, 5, MOZI_CYAN)
    
    txBox = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(4))
    tf = txBox.text_frame
    tf.text = "• 核心基因：自我管理 (Self-Management)\n• 思维底层：第一性原理 (First Principles)\n• 目标群体：具备极客潜质的新初一准学者"
    for p in tf.paragraphs:
        p.font.color.rgb = MOZI_WHITE
        p.font.size = PptPt(24)

    # Slide 3: Funnel Layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "活动转化漏斗：从人流到生源")
    
    steps = ["人流引入 (机器人)", "深度互动 (电阻测试)", "留资转化 (样品领取)", "深度咨询 (佳宁团队)"]
    for i, step in enumerate(steps):
        b = add_color_block(slide, 1 + i*2.1, 3, 2, 1, MOZI_CYAN if i < 3 else MOZI_YELLOW)
        b.text_frame.text = step
        b.text_frame.paragraphs[0].font.size = PptPt(14)
        b.text_frame.paragraphs[0].font.bold = True
        if i < len(steps)-1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, PptInches(3 + i*2.1), PptInches(3.3), PptInches(0.5), PptInches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MOZI_GRAY

    # Slide 4: Project 1 - Unitree
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "智动前沿：众擎机器人表演")
    txBox = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(4))
    tf = txBox.text_frame
    tf.text = "• 表演脚本：动态避障、台阶攀爬、语音应答互动\n• 视觉焦点：双机巡游方阵，制造震撼入场感\n• 传播目标：引发家长朋友圈高频转发，建立品牌势能"
    for p in tf.paragraphs:
        p.font.color.rgb = MOZI_WHITE
        p.font.size = PptPt(20)

    # Slide 5: Project 2 - Carbon-X
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "原子觉醒：石墨烯电阻云图")
    txBox = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(4))
    tf = txBox.text_frame
    tf.text = "• 实验内核：材料科学入门 L01\n• 互动细节：学生亲自动手测绘，理解碳材料导电特性\n• 转化点：测试完毕扫码领取“石墨烯样品礼包”"
    for p in tf.paragraphs:
        p.font.color.rgb = MOZI_WHITE
        p.font.size = PptPt(20)

    # Slide 6: Project 3 - Future Course
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "未来预告：火星/月球基地探索", MOZI_YELLOW)
    txBox = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(4))
    tf = txBox.text_frame
    tf.text = "• 课程伏笔：展示地外采矿及生存系统教具\n• 目标导向：吸引对“深空探索”感兴趣的潜在学员\n• 长期价值：建立对暑期及秋季课程的强烈心理预期"
    for p in tf.paragraphs:
        p.font.color.rgb = MOZI_WHITE
        p.font.size = PptPt(20)

    # Slide 7: Team & Responsibility
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "权责清单：协同与执行")
    txBox = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(4))
    tf = txBox.text_frame
    tf.text = "• 教学部：苏格拉底式提问引导，维护实验专业度\n• 咨询部：精准捕获意向家长，完善学员档案库\n• 工程部：电力、网络及桁架结构的绝对安全保障"
    for p in tf.paragraphs:
        p.font.color.rgb = MOZI_WHITE
        p.font.size = PptPt(20)

    # Slide 8: Safety & Risk
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "风险管控与应急响应")
    txBox = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(4))
    tf = txBox.text_frame
    tf.text = "• 硬件故障：备用电池及备用测试仪 100% 就绪\n• 极端客流：设置互动排队线，错峰进行机器人表演\n• 安全红线：机器人与观众保持 1.5 米安全距离"
    for p in tf.paragraphs:
        p.font.color.rgb = MOZI_WHITE
        p.font.size = PptPt(20)

    prs.save('/Users/zhoulin/Desktop/517天府七中墨子实验室策划方案_终极详细版.pptx')

if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    print("Ultra-Detailed Branded files generated on Desktop.")
