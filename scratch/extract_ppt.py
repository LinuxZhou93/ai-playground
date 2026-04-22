from pptx import Presentation
import os

path = "/Users/zhoulin/Library/Containers/com.tencent.xinWeChat/Data/Documents/app_data/xwechat_files/wxid_nidnlmvp7ikd22_01ad/msg/file/2026-04/2026年关于AMC8升学讲座.pptx"

if not os.path.exists(path):
    print(f"File not found: {path}")
else:
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
