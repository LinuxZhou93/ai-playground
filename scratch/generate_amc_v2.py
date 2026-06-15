from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_slide(prs, title, content_list, is_data=False):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(13, 17, 23) # Deepest Black-Blue

    title_shape = slide.shapes.title
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 255, 157) # Matrix Green
        paragraph.alignment = PP_ALIGN.LEFT

    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = ""
    for item in content_list:
        p = tf.add_paragraph()
        p.text = str(item)
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(230, 237, 243) # GitHub Style White
        p.space_after = Pt(8)
        
    if is_data:
        # Add a subtle highlight box if it's a data slide
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.05)
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 255, 157)

def create_amc_v2():
    prs = Presentation()
    
    # Slide 1: Professional Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(13, 17, 23)
    
    title = slide.shapes.title
    title.text = "2026 AMC 8 专家级宣讲：\n全球数学视野与成都温江升学实战"
    subtitle = slide.shapes.placeholders[1]
    subtitle.text = "主讲：墨子工坊首席竞赛专家 (Senior AMC Coach)\n汇集 MAA、Po-Shen Loh 与 Think Academy 全球调研精华"

    content_map = [
        ("核心议程 / AGENDA", ["• 2025-2026 全球竞赛白皮书揭秘", "• 分数线飙升：23 分时代的选拔逻辑", "• 成都名校(七中/嘉祥)的 AMC 隐形权重", "• 核心模块深度解析：代数/几何/数论/组合", "• 8 个月高效备考 roadmap 与简历策略"]),
        
        ("1. 全球视野：为什么是 AMC？", ["• 后奥数时代：逻辑深度 vs. 技巧偏门", "• 国际标准：MIT/CMU/Ivy 等名校唯一认可的数学能力背书", "• 未来人才画像：解决压力下从未接触过的陌生难题"]),
        
        ("2. 惊艳数据：2025 分数线大曝光", ["• 全球卓越奖 (DHR 1%)：23 分 (满分 25)", "• 全球优秀奖 (HR 5%)：18-19 分", "• 趋势：平均得分逐年走低，但头部竞争激烈程度翻倍", "• 结论：必须追求近乎满分的准确率"], True),
        
        ("3. 罗博深观点：数学思维的本质", ["• 'Don't teach them tricks, teach them logic.'", "• AMC 8 不是考计算，而是考查‘数学语言的翻译能力’", "• 顶级竞争者的特征：在第 21-25 题的脑力极限下依然冷静"]),
        
        ("4. 成都名校实战：简历中的“硬通货”", ["• 七中 (水井坊/天环)：数学逻辑是入围面试的门槛", "• 嘉祥锦江：AMC 1% 是顶流班型的标准配置", "• 认可度：AMC 已成为成外、实外等名校自主招生的重要参考"]),
        
        ("5. 温江家长必读：如何实现跨区逆袭", ["• 教育高地温江：嘉祥温江、七中万达等学校对奥数替代品的渴望", "• 跨区利器：若目标是主城区三校七类，AMC 是最通用的简历语言", "• 校内优势：AMC 训练能直接提升校内数学压轴题的表现"]),
        
        ("6. 考试结构解析：时间的艺术", ["• 40 分钟 / 25 题：平均每题 96 秒", "• 阅读量：每题平均翻倍至 60-80 词", "• 形式：中英双语，但高分选手必须掌握数学逻辑词汇"]),
        
        ("7. 2026 三大新趋势：深度预警", ["• 趋势一：跨领域建模 (例如在概率题中考察几何面积)", "• 趋势二：伪应用题增多 (题干冗长，干扰信息繁多)", "• 趋势三：回归线下严格监管，真实能力的唯一自证"]),
        
        ("8. 模块占比分析：分数的分布", ["• 应用题 (Word Problems)：30% (区分度的核心)", "• 代数 (Algebra)：25% (基本盘，不能丢分)", "• 几何 (Geometry)：20% (直觉与辅助线)", "• 数论与组合 (NT & Counting)：25% (冲刺难题的关键)"], True),
        
        ("9. 核心模块：代数部分的进化", ["• 从简单的方程到复杂的数列规律与比率分析", "• 考点：行程问题、工程问题、时钟问题", "• 专家建议：掌握逻辑消元法，减少笔算步骤"]),
        
        ("10. 核心模块：几何部分的“辅助线”", ["• 考点：相似三角形、圆面积、切割法、勾股定理", "• 难点：三维空间的二维展开与体积计算", "• 技巧：几何问题的本质是寻找‘不变的关联’"]),
        
        ("11. 核心模块：数论——中国学生的盲区", ["• 考点：整除特征、质因数分解应用、余数系统、进制转换", "• 困境：体制内涉及极少，AMC 却作为高分分水岭", "• 策略：系统构建数论模型，建立对数字的直觉"]),
        
        ("12. 核心模块：计数与概率——脑力极限", ["• 考点：组合数学、概率期望、维恩图、最值逻辑", "• 挑战：分类讨论的严谨性，不重不漏", "• 案例：如何通过逻辑排除掉 50% 的干扰选项"]),
        
        ("13. 试卷分布攻略：梯队战术", ["• 1-10 题 (Green Zone)：5-8 分钟横扫，准确率 100%", "• 11-20 题 (Yellow Zone)：15-20 分钟，攻坚逻辑陷阱", "• 21-25 题 (Red Zone)：10-15 分钟，专家选手的封神时刻"]),
        
        ("14. 专家解题神技 (一)：排除与代入", ["• 既然是选择题，逻辑排除比暴力计算更有效", "• 特殊值法：将抽象字母具象化，瞬间破题"]),
        
        ("15. 专家解题神技 (二)：图形可视化", ["• 将枯燥的应用题转化为几何图形或坐标轴", "• 一图胜千言：逻辑可视化的魅力"]),
        
        ("16. 备考阶段一：知识扫盲 (1-3月)", ["• 系统学习 Pre-algebra 与 Algebra 基础", "• 确保没有知识死角，解决‘看不懂题’的问题"]),
        
        ("17. 备考阶段二：专题攻坚 (4-6月)", ["• 针对数论、组合、复杂几何进行专项训练", "• 建立每个模块的解题‘模板库’"]),
        
        ("18. 备考阶段三：真题磨炼 (7-9月)", ["• 近 15 年真题地毯式过关", "• 重点在于：找到出题人的思维逻辑陷阱"]),
        
        ("19. 备考阶段四：极限冲刺 (10-1月)", ["• 全真 40 分钟闭卷模拟", "• 重点在于：在倒计时红色预警下的抗压与决策"]),
        
        ("20. G3-G4：播种兴趣的黄金季", ["• 目标：袋鼠竞赛 (Math Kangaroo) 衔接", "• 重点：培养对数字的热爱，而非刷题压力"]),
        
        ("21. G5-G6：冲击 1% 的巅峰期", ["• 目标：DHR 全球卓越奖", "• 现状：这是进入名校简历库含金量最高的时期"]),
        
        ("22. G7-G8：最后的反向加码", ["• 目标：校内分班锁定或直接跨段备战 AMC 10", "• 意义：证明即使在课业繁重的初中，数学思维依然拔尖"]),
        
        ("23. 简历展示的艺术：如何写才高大上", ["• 错误写法：AMC 8 成绩 19 分", "• 专家建议：'2025 AMC 8 全球优秀奖 (Top 5%)，数论模块满分'", "• 突出百分比，而非具体分数"]),
        
        ("24. 给家长的心态建设：是陪伴而非督促", ["• 竞赛是试金石，发现孩子的天赋区域", "• 哪怕没获奖，逻辑思维的提升对中考数学是全方位的降维打击"]),
        
        ("25. 为什么很多“尖子生”会兵败 AMC？", ["• 误区一：习惯了题海战术，无法应对即兴解题", "• 误区二：时间分配不规范，在简单的第 5 题卡死 10 分钟", "• 专家贴士：专注‘提问的艺术’"]),
        
        ("26. 墨子工坊的特色：AI 赋能竞赛", ["• AI 定制错题分账：只练孩子不会的那 5% 的漏洞", "• 可视化逻辑解析：让抽象的数论变成生动的动画"]),
        
        ("27. 教练团队：深耕成都教育圈", ["• 我们不仅懂 AMC，我们更懂成都名校的招生偏好", "• 温江家长的‘升学内参’"]),
        
        ("28. 成功案例分享：从温江到七中", ["• 案例：某温江普小学生，通过 1 年系统训练斩获 DHR，获主城区顶级学校垂青", "• 核心：提前 18 个月布局的力量"]),
        
        ("29. 资源整合：推荐书单与网站", ["• 书籍：AoPS (Art of Problem Solving) 全系列", "• 工具：Alcumus 在线练习平台", "• 墨子工坊专属真题精析包"]),
        
        ("30. 2026 考试时间表预演", ["• 2025.09：大纲下载与初测", "• 2025.11：报名确认", "• 2026.01.23：决战 40 分钟"]),
        
        ("31. 常见问题 FAQ (一)：英语水平要求", ["• 实际上考查的是‘数学英语’", "• 掌握 100 个核心词汇即可无障碍通关"]),
        
        ("32. 常见问题 FAQ (二)：要不要学奥数？", ["• 奥数是‘深度’，AMC 是‘广度+速度’", "• 两者相辅相成，AMC 是检验奥数成果的最佳国际化平台"]),
        
        ("33. 常见问题 FAQ (三)：费用与考点", ["• 详细列出报名费与授权考点分布（包括成都本地地址）"]),
        
        ("34. 专家视点：关于 AIME 的提前布局", ["• 既然开始卷，就卷到底：G6 冲击 AMC 8，G7 直接挑战 AMC 10 晋级 AIME", "• 这种‘跨级’履历在名校眼中是极具杀伤力的"]),
        
        ("35. 文江本地教育趋势：温江七中的逻辑", ["• 温七创新班分班考中的竞赛题型调研分享"]),
        
        ("36. 数学，不仅是学科，更是权力", ["• 解析逻辑思维如何帮助孩子在未来 AI 时代不被替代"]),
        
        ("37. 墨子工坊课程导引：模块课 v.s. 长期班", ["• 针对不同基础的孩子提供定制化的陪跑方案"]),
        
        ("38. 限时福利：现场测评报告", ["• 只要是今晚到场的家长，均可获得一次免费的 AMC 潜力测试"]),
        
        ("39. 结语：给孩子一个证明自己的舞台", ["• 数学竞赛不应是累赘，而应是高光时刻", "• 愿每个孩子都能在逻辑的海洋中自信起航"]),
        
        ("40. 结束页：添加专家微信，获取全套资料", ["• 扫描二维码：领取《2026 AMC 8 备考白皮书》", "• 感谢聆听，今晚愿与各位家长共同守护孩子的逻辑之光"])
    ]

    for i, item in enumerate(content_map):
        title = item[0]
        content = item[1]
        is_data_slide = item[2] if len(item) > 2 else (i%5==0 and i>0)
        add_slide(prs, title, content, is_data=is_data_slide)
    
    path = "/Users/zhoulin/Desktop/2026_AMC_Elite_Expert_v2.pptx"
    prs.save(path)
    print(f"PPT generated at: {path}")

create_amc_v2()
