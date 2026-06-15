import docx
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- MOZI BRANDING ---
MOZI_DARK = PptRGBColor(7, 8, 12)
MOZI_CYAN = PptRGBColor(0, 240, 255)
MOZI_YELLOW = PptRGBColor(255, 214, 0)
MOZI_WHITE = PptRGBColor(245, 245, 247)
MOZI_DEEP_BLUE = PptRGBColor(15, 23, 42)

def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_manual_title(slide, text, color=MOZI_CYAN):
    txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(9), PptInches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = color
    p.font.bold = True
    p.font.size = PptPt(36)

def add_color_block(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_full_plan_content(doc):
    t = doc.add_heading('【5.17天府七中】墨子实验室校园开放日·全场景全流程深度执行全案', 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_heading('一、 品牌战略与教育愿景', level=1)
    doc.add_paragraph('墨子实验室（Mozi Lab）作为全栈科技特长生培养基地，不仅是硬核科技的展示场，更是未来科创人才的“第一战场”。本次活动的总体方针是：利用极致的视觉冲击力（众擎机器人）和深度的科学沉究（Carbon-X）在准初一群体中建立绝对的技术护城河。')
    doc.add_heading('二、 现场布防与视觉识别系统', level=1)
    doc.add_paragraph('桁架广告：6m x 3.5m 超大喷绘，背景采用深邃星空结合科技蓝光效，突出“墨子学者”身份标识。\n区域布局：\n1. 表演区：机器人动态巡演，配备警戒线与导流围栏。\n2. 实验区：4组实验台，配备数字示波器、万用表及石墨烯测试样品。\n3. 咨询区：佳宁老师团队坐镇，配备iPad学员信息录入系统及纸质手册。')
    doc.add_heading('三、 详细互动脚本与执行细节', level=1)
    doc.add_heading('3.1 众擎仿生机器人：全能极客表演', level=2)
    doc.add_paragraph('• 表演周期：每30分钟进行一次。表演曲目包括：原地平衡、障碍物跨越、模拟交互握手。\n• 互动Q&A：导员向学生提问：“你知道机器人是怎么感知到前面的障碍物的吗？”引出超声波与激光雷达知识点。\n• 转化动作：合影完毕后，告知学生“去实验区领一份机器人同款石墨烯材料”。')
    doc.add_heading('3.2 Carbon-X：石墨烯精英班核心实验', level=2)
    doc.add_paragraph('• 实验步骤：每组5分钟。通过多点测试，学生会发现石墨烯纸各处的阻值并不完全相同。以此解释材料的微观结构与宏观性能的关系。\n• 伏笔铺设：老师在实验结束时说：“这只是石墨烯入门。在我们的精英学者营，我们将用这些材料直接制作传感器，装在刚才那个机器人身上。”')
    doc.save('/Users/zhoulin/Desktop/517天府七中开放日全案_全量细节版.docx')

def generate_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_color_block(slide, 0, 0, 10, 0.2, MOZI_CYAN)
    add_color_block(slide, 0, 0.2, 0.2, 7.5, MOZI_CYAN)
    tx = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(2))
    p = tx.text_frame.paragraphs[0]
    p.text = "5.17 天府七中开放日"
    p.font.size = PptPt(64)
    p.font.bold = True
    p.font.color.rgb = MOZI_CYAN
    p2 = tx.text_frame.add_paragraph()
    p2.text = "墨子实验室：深度引流与展示全案"
    p2.font.size = PptPt(32)
    p2.font.color.rgb = MOZI_WHITE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "01 | 品牌核心：科技护城河")
    add_color_block(slide, 0.5, 2, 4.2, 4, MOZI_DEEP_BLUE)
    add_color_block(slide, 5.3, 2, 4.2, 4, MOZI_DEEP_BLUE)
    tx1 = slide.shapes.add_textbox(PptInches(0.6), PptInches(2.2), PptInches(4), PptInches(3))
    tf1 = tx1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "【核心使命】\n培养具备物理直觉与\n工程能力的科创领袖。\n六维评估体系全面赋能。"
    p1.font.size = PptPt(18); p1.font.color.rgb = MOZI_WHITE
    tx2 = slide.shapes.add_textbox(PptInches(5.4), PptInches(2.2), PptInches(4), PptInches(3))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "【转化逻辑】\n感官冲击(A) -> \n深度探究(B) -> \n礼包钩子(C) -> \n线索采集(D)"
    p2.font.size = PptPt(18); p2.font.color.rgb = MOZI_WHITE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "02 | 转化漏斗：流量到线索")
    for i in range(4):
        w = 8 - i*1.5; x = (10 - w)/2; y = 2 + i*1.1
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(0.8))
        shape.fill.solid(); shape.fill.fore_color.rgb = MOZI_CYAN
        shape.text_frame.text = ["机器人触达", "石墨烯互动", "礼包扫码", "咨询转化"][i]
        shape.text_frame.paragraphs[0].font.size = PptPt(16); shape.text_frame.paragraphs[0].font.color.rgb = MOZI_DARK

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "03 | 执行：众擎巡游脚本")
    add_color_block(slide, 0.5, 1.5, 9, 5, MOZI_DEEP_BLUE)
    tx = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.8), PptInches(8.5), PptInches(4.5))
    tf = tx.text_frame
    tf.text = "• 极客表演：每逢整点 15 分钟挑战赛\n• 视觉焦点：双机巡游阵列，制造震撼\n• 互动：设置“未来指挥官”合影区域\n• 后勤：双组电池热切换保障 8 小时巡航"
    for para in tf.paragraphs: para.font.size = PptPt(22); para.font.color.rgb = MOZI_WHITE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, MOZI_DARK)
    add_manual_title(slide, "04 | 伏笔：火星/月球基地", MOZI_YELLOW)
    add_color_block(slide, 0.5, 1.5, 9, 5, MOZI_DEEP_BLUE)
    tx = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.8), PptInches(8.5), PptInches(4.5))
    tf = tx.text_frame
    tf.text = "• 内容：地外采矿教具与密封仓展示\n• 策略：不仅仅是模型，而是挑战计划\n• 转化：仅开放 50 个内测席位预约\n• 目标：将今日流量锁定为秋季核心生源"
    for para in tf.paragraphs: para.font.size = PptPt(22); para.font.color.rgb = MOZI_WHITE

    prs.save('/Users/zhoulin/Desktop/517天府七中开放日全案_全量细节版.pptx')

if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    print("Full-Detail Branded files generated on Desktop.")
