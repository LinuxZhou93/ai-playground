import docx
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- MCKINSEY CLEAN STYLE (GOOGLE-ISH BUT MORE CORPORATE) ---
M_NAVY = PptRGBColor(5, 28, 72)
M_BLUE = PptRGBColor(0, 115, 207)
M_WHITE = PptRGBColor(255, 255, 255)
M_GRAY = PptRGBColor(242, 242, 242)
M_TEXT = PptRGBColor(33, 33, 33)

def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_mckinsey_title(slide, text, color=M_NAVY):
    txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.4), PptInches(9), PptInches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = color
    p.font.bold = True
    p.font.size = PptPt(28)

def add_full_mckinsey_plan(doc):
    t = doc.add_heading('【5.17天府七中】墨子实验室开放日战略策划方案 (McKinsey Framework)', 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_heading('1. 执行摘要 (Executive Summary)', level=1)
    doc.add_paragraph('本项目旨在利用天府七中5.17校园开放日产生的2000-3000名高价值流量，通过“高冲击力视觉、深参与感互动、长链路心理锚点”三大战略引擎，确立墨子实验室在区域科创教育市场的领导地位，并实现高净值潜客的全量录入。')
    
    doc.add_heading('2. 战略目标与KPI (Strategic Intent & KPIs)', level=1)
    doc.add_paragraph('• 品牌定位：确立“全栈科技特长生培养”第一品牌形象。\n• 转化目标：实现30%的到场家庭有效留资（约600-900条线索）。\n• 获客成本：利用现有硬件及复用物料，将单条线索获取成本控制在5元以下。')
    
    doc.add_heading('3. 受众洞察与价值主张 (Audience Insights & Value Proposition)', level=1)
    doc.add_paragraph('天七家长群体特征：教育焦虑较高，追求差异化背景，对名校科创路径敏感。\n价值主张：不仅是玩机器人，而是通过材料科学与工程思维，让孩子具备进入全球Top高校的科创护城河。')
    
    doc.add_heading('4. 三位一体活动设计 (The Pillars of Execution)', level=1)
    doc.add_heading('4.1 视觉冲击岛：众擎机器人 (Sensory Disruption)', level=2)
    doc.add_paragraph('利用双机阵列展示全地形动态平衡，建立技术权威感。')
    doc.add_heading('4.2 认知沉浸岛：Carbon-X 实验 (Cognitive Engagement)', level=2)
    doc.add_paragraph('通过石墨烯电阻实测，将抽象科学转化为可量化的工程实验。')
    doc.add_heading('4.3 心理锚点岛：火星探险预热 (Psychological Seed)', level=2)
    doc.add_paragraph('发布未来月球/火星课程，建立“未完待续”的长期粘性。')
    
    doc.add_heading('5. 运营转化与SOP (Operational Funnel)', level=1)
    doc.add_paragraph('入场(机器人) -> 停留(扫码领样品) -> 深度体验(电阻测试) -> 顾问一对一(潜质评估) -> 后期私域运营。')
    
    doc.add_heading('6. 风险管控方案 (Risk & Mitigation)', level=1)
    doc.add_paragraph('详细列出设备冗余、流量过载、及极端天气下的B计划。')
    
    doc.save('/Users/zhoulin/Desktop/517天府七中策划方案_麦肯锡战略版.docx')

def generate_pptx():
    prs = Presentation()
    
    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, M_WHITE)
    add_mckinsey_title(slide, "5.17 天府七中校园开放日战略策划", M_NAVY)
    
    tx = slide.shapes.add_textbox(PptInches(0.5), PptInches(2.5), PptInches(9), PptInches(2))
    p = tx.text_frame.paragraphs[0]
    p.text = "构建“墨子学者”精英教育生态的流量闭环"
    p.font.size = PptPt(36); p.font.bold = True; p.font.color.rgb = M_BLUE
    
    p2 = tx.text_frame.add_paragraph()
    p2.text = "Strategic Proposal | May 2026"; p2.font.size = PptPt(18); p2.font.color.rgb = M_TEXT

    # Slide 2: Strategic Context
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, M_GRAY)
    add_mckinsey_title(slide, "1. 战略意图：流量获取与品牌锚定")
    tx = slide.shapes.add_textbox(PptInches(0.5), PptInches(1.5), PptInches(9), PptInches(4))
    tf = tx.text_frame
    tf.text = "• 市场背景：天府七中开放日聚集了全成都市最精准的准初一科创家庭。\n• 战略目标：通过硬核科技演示，在2小时内完成品牌认知到行动转化的全过程。\n• KPI设定：600+ 深度留资线索，20% 暑期营预约率。"
    for p in tf.paragraphs: p.font.size = PptPt(20); p.font.color.rgb = M_TEXT

    # Slide 3: Audience & Value
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, M_WHITE)
    add_mckinsey_title(slide, "2. 用户画像与核心价值主张")
    tx = slide.shapes.add_textbox(PptInches(0.5), PptInches(1.5), PptInches(9), PptInches(4))
    tf = tx.text_frame
    tf.text = "• 用户画像：高知中产家庭，关注新中考/高考趋势，对科技特长生路径有刚需。\n• 核心主张：墨子实验室不仅是教具，而是“全栈科技特长生”的闭环培养图谱。\n• 核心优势：基于北大课程内核，融合前沿材料学与仿生机器人技术。"
    for p in tf.paragraphs: p.font.size = PptPt(20); p.font.color.rgb = M_TEXT

    # Slide 4: Pillars
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, M_GRAY)
    add_mckinsey_title(slide, "3. 落地执行：三位一体展示逻辑")
    # 3 Boxes
    for i, title in enumerate(["视觉冲击岛", "认知沉浸岛", "心理锚点岛"]):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.5 + i*3.1), PptInches(2), PptInches(2.8), PptInches(4))
        shape.fill.solid(); shape.fill.fore_color.rgb = M_NAVY
        shape.text_frame.text = f"{title}\n\n{['众擎机器人演示', 'Carbon-X 实测', '月球/火星课预告'][i]}"
        shape.text_frame.paragraphs[0].font.size = PptPt(16); shape.text_frame.paragraphs[0].font.bold = True

    # Slide 5: Risk & Team
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, M_WHITE)
    add_mckinsey_title(slide, "4. 运营保障与风险管理")
    tx = slide.shapes.add_textbox(PptInches(0.5), PptInches(1.5), PptInches(9), PptInches(4))
    tf = tx.text_frame
    tf.text = "• 资源：佳宁老师团队负责闭环转化，工程部负责电力及结构安全。\n• 风险：设备冗余（众擎备机）、流量分发（排队系统）、离线留资（纸质备份）。\n• 迭代：活动当日每2小时进行一次运营复盘与话术优化。"
    for p in tf.paragraphs: p.font.size = PptPt(20); p.font.color.rgb = M_TEXT

    prs.save('/Users/zhoulin/Desktop/517天府七中策划方案_麦肯锡战略版.pptx')

if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    print("McKinsey Strategic Proposal generated on Desktop.")
