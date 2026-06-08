import docx
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- MCKINSEY CLEAN STYLE ---
M_NAVY = PptRGBColor(5, 28, 72)
M_BLUE = PptRGBColor(0, 115, 207)
M_WHITE = PptRGBColor(255, 255, 255)
M_GRAY = PptRGBColor(242, 242, 242)
M_TEXT = PptRGBColor(33, 33, 33)

def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_mckinsey_title(slide, text, color=M_NAVY):
    txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.4), PptInches(9), PptInches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.color.rgb = color
    p.font.bold = True
    p.font.size = PptPt(28)

def generate_docx():
    doc = docx.Document()
    t = doc.add_heading('【5.17天府七中】墨子实验室开放日战略策划方案 (McKinsey Framework)', 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_heading('1. 执行摘要', level=1)
    doc.add_paragraph('本项目旨在利用天府七中 5.17 校园开放日的流量高峰，通过“视觉震撼-认知深耕-转化闭环”的三位一体战略，将墨子实验室确立为区域高端科创教育的领军品牌。')
    doc.add_heading('2. 战略意图与核心 KPI', level=1)
    doc.add_paragraph('• 目标：实现 2000+ 准初一家庭的高频触达，沉淀 600+ 高价值销售线索。\n• KPI：留资率 > 30%，品牌提及率 > 50%。')
    doc.add_heading('3. 活动落地执行架构', level=1)
    doc.add_paragraph('3.1 众擎机器人：作为流量入口（Sensory Disruption），负责制造视觉焦点。\n3.2 Carbon-X 实验：作为认知锚点（Cognitive Engagement），负责建立专业护城河。\n3.3 佳宁团队：作为转化闭环（Conversion Engine），负责精准意向捕获。')
    doc.save('/Users/zhoulin/Desktop/517天府七中策划方案_麦肯锡战略版.docx')

def generate_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, M_WHITE)
    add_mckinsey_title(slide, "5.17 天府七中校园开放日战略策划")
    tx = slide.shapes.add_textbox(PptInches(0.5), PptInches(2.5), PptInches(9), PptInches(2))
    p = tx.text_frame.paragraphs[0]
    p.text = "构建“墨子学者”精英教育生态的流量闭环"
    p.font.size = PptPt(36); p.font.bold = True; p.font.color.rgb = M_BLUE
    prs.save('/Users/zhoulin/Desktop/517天府七中策划方案_麦肯锡战略版.pptx')

if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    print("McKinsey Strategic Proposal generated on Desktop.")
