import docx
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.text import PP_ALIGN

def generate_docx():
    doc = docx.Document()
    
    # Title
    title = doc.add_heading('天府七中墨子实验室校园开放日（5.17）展示方案 · 优化版', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Background
    doc.add_heading('一、 活动核心背景', level=1)
    p = doc.add_paragraph('本次活动针对2000-3000名新初一准学生及家长。展示不仅是科技实力的博览，更是招生引流的“黄金窗口”。方案重点强化“伴手礼引流”与“课程伏笔”机制。')
    
    # Location
    doc.add_heading('二、 点位与视觉呈现', level=1)
    p = doc.add_paragraph()
    p.add_run('核心坐标：').bold = True
    p.add_run('报告厅正门口（全校流量核心）。')
    
    p = doc.add_paragraph()
    p.add_run('视觉物料：').bold = True
    p.add_run('6m x 3m 科技感桁架广告牌，主题：“碳基极客，智绘未来”。大屏幕循环播放墨子实验室生态宣传片。')
    
    # Projects
    doc.add_heading('三、 互动展示项目（深度优化）', level=1)
    
    items = [
        ("1. 众擎机器人：全地形巡游表演", "机器人作为门面担当，在报告厅门口进行避障、招呼及舞蹈表演。现场设置“机器人合影区”，引导家长拍照分享朋友圈。"),
        ("2. 石墨烯“伴手礼”引流区（核心Hook）", "每位到场学生扫码关注公众号/添加老师微信，即可领取一份“石墨烯导电薄膜”样品及科普卡片。在卡片中埋下“实验室深度探秘”的伏笔。"),
        ("3. 火星/月球基地：未来课程预告", "替代原有的空间站模型，展示“火星基地探索”或“月球资源开发”课程的先导课件与实物教具，吸引学生对暑期/秋季课程的期待。"),
        ("4. 石墨烯电阻云图实测", "学生现场操作万用表，绘制自己的“电阻云图”，通过直观的数据变化理解新材料的科学底层逻辑。"),
        ("5. 碳基极客工坊：DIY 导电电路", "使用导电笔和石墨烯纸现场制作能发光的科技明信片，强调“材料+创意”的结合。")
    ]
    
    for title_text, desc in items:
        p = doc.add_paragraph()
        p.add_run(title_text).bold = True
        p.add_run(f'\n{desc}')
    
    # Workflow
    doc.add_heading('四、 流程与引流逻辑', level=1)
    doc.add_paragraph('08:30 - 09:00：桁架搭建与机器人热身\n09:00 - 16:30：滚动展示。每逢整点进行机器人集中表演。\n引流逻辑：扫码领材料 -> 互动体验 -> 咨询登记 -> 预约后续实验室开放日。')
    
    # Cost
    doc.add_heading('五、 成本控制建议', level=1)
    doc.add_paragraph('• 众擎机器人为已有设备，无需新增租赁费。\n• 石墨烯样品及导电笔耗材建议提前集采，单人综合成本控制在5元以内。\n• 佳宁老师团队负责现场引导与信息录入，实现招生转化。')

    doc.save('/Users/zhoulin/Desktop/517天府七中墨子实验室开放日策划方案_优化版.docx')

def generate_pptx():
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "5.17 天府七中开放日展示方案"
    slide.placeholders[1].text = "碳基极客 · 智绘未来\n墨子实验室团队 · 深度优化版"
    
    # Slide 2: Pain Points
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "核心策略：如何让2000+学生留下来？"
    content = slide.placeholders[1]
    content.text = "• 强视觉：大型行架 + 众擎机器人巡游\n• 强关联：人手一份石墨烯材料样品（扫码领）\n• 强期待：火星/月球基地探索课程“伏笔”"
    
    # Slide 3: Layout
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "点位安排：报告厅门口 C 位"
    content = slide.placeholders[1]
    content.text = "• 6m x 3m 桁架，打造沉浸式实验室入口感\n• 设置拍照打卡点，引发社交裂变\n• 专设家长咨询区（佳宁老师团队配合）"
    
    # Slide 4: Highlights 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "硬核互动 1：众擎机器人 & 石墨烯礼包"
    content = slide.placeholders[1]
    content.text = "• 动态表演吸引人流\n• 扫码领样品实现拉新\n• 卡片留伏笔：引导预约实验室深度参观"
    
    # Slide 5: Highlights 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "硬核互动 2：探索课程 & 极客工坊"
    content = slide.placeholders[1]
    content.text = "• 火星/月球基地课程先导展示\n• 电阻云图实测（材料科学深度认知）\n• DIY 导电电路（带得走的科技作品）"
    
    # Slide 6: Operation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "执行与成本控制"
    content = slide.placeholders[1]
    content.text = "• 耗材人均成本 < 5元\n• 复用实验室已有高新设备\n• 佳宁老师团队负责精准意向挖掘"

    prs.save('/Users/zhoulin/Desktop/517天府七中墨子实验室开放日策划方案_优化版.pptx')

if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    print("Optimized files generated on Desktop.")
